import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE
from PIL import Image
import csv


class PPTLoader:
    def __init__(self, file_path, output_dir):
        self.file_path = file_path
        self.output_dir = output_dir
        self.presentation = Presentation(file_path)

    def extract_text(self):
        # Extract text from all slides
        text = ""
        for slide_num, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        # Save the extracted text to a file
        text_path = os.path.join(self.output_dir, 'text', 'extracted_text.txt')
        with open(text_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"Text extracted and saved at: {text_path}")

    def extract_images(self):
        images = []
        image_folder = os.path.join(self.output_dir, 'images')

        for slide_num, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    img_ext = image.ext
                    img_path = os.path.join(image_folder, f'image_slide_{slide_num + 1}.{img_ext}')

                    with open(img_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    images.append(img_path)
        print(f"Images extracted and saved at: {image_folder}")
        return images

    def extract_tables(self):
        tables = []
        table_folder = os.path.join(self.output_dir, 'tables')

        for slide_num, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)

                    tables.append(table_data)

                    # Save table data to CSV
                    table_csv_path = os.path.join(table_folder, f"table_slide_{slide_num + 1}.csv")
                    with open(table_csv_path, "w", encoding="utf-8") as f:
                        writer = csv.writer(f)
                        writer.writerows(table_data)
                    print(f"Table from slide {slide_num + 1} saved at: {table_csv_path}")
        return tables

    def extract_links(self):
        links = []
        link_folder = os.path.join(self.output_dir, 'links')
        link_csv_path = os.path.join(link_folder, "extracted_links.csv")

        for slide_num, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                links.append({
                                    "slide": slide_num + 1,
                                    "url": run.hyperlink.address
                                })

        # Save links in the "links" subfolder
        with open(link_csv_path, "w", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Slide", "URL"])
            for link in links:
                writer.writerow([link["slide"], link["url"]])
        print(f"Links extracted and saved at: {link_csv_path}")
        return links
